"""Shared pytest fixtures and helpers for mindvault tests."""

from __future__ import annotations

import shutil
from pathlib import Path

import pytest


FIXTURES_DIR = Path(__file__).parent / "fixtures"


@pytest.fixture
def obsidian_vault_fixture(tmp_path: Path) -> Path:
    """Copy the static Obsidian vault fixture into a temp dir so tests can mutate it."""
    src = FIXTURES_DIR / "obsidian_vault"
    dst = tmp_path / "vault"
    shutil.copytree(src, dst)
    return dst


@pytest.fixture
def office_docs_dir(tmp_path: Path) -> Path:
    """Generate docx / xlsx / pptx fixtures on the fly into a temp dir.

    Keeping these as generated-at-test-time (instead of committed binaries) avoids
    binary files in the repo and exercises the python-docx / openpyxl / python-pptx
    dependency chain as a side effect.
    """
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation

    office_dir = tmp_path / "office"
    office_dir.mkdir()

    # docx
    doc = Document()
    doc.add_heading("프로젝트 A 제안서", 0)
    doc.add_paragraph("MindVault 문서 추출 테스트입니다.")
    doc.add_paragraph("핵심 기능: 그래프 + 위키 + 검색.")
    doc.save(str(office_dir / "proposal.docx"))

    # xlsx
    wb = Workbook()
    ws = wb.active
    ws.title = "매출"
    ws.append(["월", "매출", "비용"])
    ws.append(["1월", 1000000, 400000])
    ws.append(["2월", 1200000, 500000])
    wb.save(str(office_dir / "sales.xlsx"))

    # pptx
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "2026 로드맵"
    s1.placeholders[1].text = "MindVault 확장 계획"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "주요 기능"
    prs.save(str(office_dir / "roadmap.pptx"))

    return office_dir


@pytest.fixture
def clean_mindvault_config(tmp_path: Path, monkeypatch):
    """Redirect ~/.mindvault to a temp dir so tests don't clobber real user config."""
    fake_home = tmp_path / "home"
    fake_home.mkdir()
    monkeypatch.setenv("HOME", str(fake_home))
    # mindvault.config uses Path.home() which reads HOME at call time
    yield fake_home
